"""
PPTX CHARTER KT THEME V1.5

Developed by: Fernando Monrroy fernando.monrroy@kantar.com

This custom library is a wrapper for Python-PPTX library. 
It has custom tables and charter functions that generate PPTX objects already alligned with the current Kantar Theme (as of Auguts 2025). 
"""

# Initialization 
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE, XL_AXIS_CROSSES
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor 
import lxml.etree as ET
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# Standard color pallete
standard_pallete = [
    RGBColor(0,96,255), #kantar_blue_strong
    RGBColor(128,42,183), #kantar_purple_strong
    RGBColor(0,229,186), #kantar_teal_strong
    RGBColor(0,182,0), #kantar_green_strong
    RGBColor(250,0,40), #kantar_red_strong
    RGBColor(254,219,0), #kantar_yellow_strong
    RGBColor(255,80,0), #kantar_orange_strong
    RGBColor(181,115,223), #kantar_blue_mid
    RGBColor(101,159,255), #kantar_purple_mid
    RGBColor(86,255,223), #kantar_teal_mid
    RGBColor(58,255,58), #kantar_green_mid
    RGBColor(255,99,123), #kantar_red_mid
    RGBColor(254,233,101), #kantar_yellow_mid
    RGBColor(255,149,101), #kantar_orange_mid
    RGBColor(230,208,244), #kantar_blue_light
    RGBColor(204,223,255), #kantar_purple_light
    RGBColor(198,254,244), #kantar_teal_light
    RGBColor(189,255,189), #kantar_green_light
    RGBColor(255,203,211), #kantar_red_light
    RGBColor(255,247,203), #kantar_yellow_light
    RGBColor(255,220,204), #kantar_orange_light
]


# Slide object mapping functions
def get_slide_placeholders(slide_object):
    """
    This is a function that prints a list of the placeholdes a slide object has. Useful for mapping the available elements for edition. 
    """
    for placeholder in slide_object.placeholders:
        print(f"{placeholder.placeholder_format.idx}: {placeholder.name}")


def get_slide_shapes(slide_object):
    """
    This is a function that prints a list of the shapes a slide object has. Useful for mapping the available elements for edition. 
    """
    for shape in slide_object.shapes:
        print(f'{shape.shape_id}:{shape.name}')


def get_presentation_layouts(presentation_object):
    """
    Function that prints a list of the layouts of a presentation object. 
    """
    counter = 0
    for layout in presentation_object.slide_layouts:
        print(f'{counter}: {layout.name}')
        counter += 1


# Auxiliary function that checks if all columns for a multilevel index are str type 
def check_columns_strtype(df=pd.DataFrame(), all_columns=[]):
    
    """
    Function that checks if the columns in the all_columns list are string type, otherwise it raises an error listing all the non str type columns. 
    """

    non_string_columns = {}

    for col in all_columns:
        if not pd.api.types.is_string_dtype(df[col]):
            non_string_columns[col] = df[col].dtype

    if non_string_columns:
        error_message = "The following columns are not of string type:\n"
        for col, dtype in non_string_columns.items():
            error_message += f"{col}: {dtype}\n"
        raise ValueError(error_message)


# Multilevel crosstab function
def get_crosstab(
        stacked_df=pd.DataFrame(), 
        index_vars=[], 
        column_vars=[], 
        value_var='', 
        filter_mask=None, 
        agg_function='sum', 
        normalize: "bool | Literal[0, 1, 'all', 'index', 'columns']" = False, 
        str_separator=' '
        ):
    
    """
    Extends the functionality of the pd.crosstab function. 
    It allows to prefilter the data for the crosstab with the filter mask parameter, this parameter needs a masking pd.Series that will apply as filter. 
    Also it allows to nest variables for the index and columns elements of the crosstab. 
    
    Note: for this to function the index_vars and column_vars need to be of string type. Make sure to convert them beforehand into a string format. 
    """

    # Applies filter (if applicable)
    if filter_mask is None:
        _filtered_df = stacked_df.copy()
    else:
        _filtered_df = stacked_df[filter_mask].copy()


    # Checks if all the columns are string type, othewise it raises an error
    _all_columns = []
    _all_columns.extend(column_vars)
    _all_columns.extend(index_vars)

    check_columns_strtype(_filtered_df, _all_columns)


    # Creates nested index and column variables 
    _filtered_df['agg_index_var'] = _filtered_df[index_vars].agg(str_separator.join, axis=1)
    _filtered_df['agg_column_var'] = _filtered_df[column_vars].agg(str_separator.join, axis=1)


    # Generates the crosstab with the nested index and column variables
    _crosstab = pd.crosstab(
        index=_filtered_df['agg_index_var'],
        columns=_filtered_df['agg_column_var'],
        values=_filtered_df[value_var],
        aggfunc=agg_function,
        normalize=normalize
    )

    return _crosstab


# Chart maker function 
def chart_maker(
        press_object:Presentation,
        layout_object,
        use_slide_object:bool=False,
        slide_object=None,
        slide_title:str='title', 
        stacked_df:pd.DataFrame=pd.DataFrame(), 
        value_var:str='', 
        category_vars:list=[], 
        series_vars:list=[], 
        iterable_vars:list=[], 
        grouping_vars:list=[],
        filter_mask = None, 
        chart_type:XL_CHART_TYPE=XL_CHART_TYPE.COLUMN_STACKED,
        show_legend:bool=True,
        show_category_axis:bool=True,
        show_data_labels_perc:bool=False,
        show_data_labels_value:XL_LABEL_POSITION|None=None,
        pos_horizontal:Inches=Inches(0.39),
        pos_vertical:Inches=Inches(1.87),
        width:Inches=Inches(12.54),
        height:Inches=Inches(4.37),
        agg_function='sum', 
        color_pallete:list=standard_pallete,
        number_format_series:str='#,##0',
        zero_axis_start:bool=True,
        crosstab_override:pd.DataFrame=None
        ):
    
    """
    This function generates graphs slides in 3 modes of operations: 

    - Single graph in a slide -
    This is the most simple mode, this will create a new slide with a single graph on it. For this mode of operation both iterable_vars and grouping vars should be left empty. 

    - Multiple single graphs slides -
    This mode generates multiple new slides with a single graph on each. The graph will be created by looping on the unique elements of iterable_vars object. 

    - Multiple graph triplet slides - 
    This is the most complex mode, this will create up to 3 graphs in a single slide per group of every unique element in the grouping_vars object. 
    If a group would create more than 3 graphs it will create more slides as triplets of graphs.
    This mode require sboth iterable_vars and grouping_vars lists. 

    """

    # Filter the data
    if filter_mask is None:
        _filtered_df = stacked_df.copy()
    else:
        _filtered_df = stacked_df[filter_mask].copy()

    if (len(iterable_vars) == 0) & (len(grouping_vars) == 0):

        # Prepare the data

        # If there is no override, it'll use the specified parameters in the function to make the crosstab
        if crosstab_override is None:
            _crosstab = get_crosstab(
                stacked_df=_filtered_df,
                index_vars=category_vars,
                column_vars=series_vars,
                value_var=value_var,
                agg_function=agg_function
            ).fillna(0)
        
        # If there is a crosstab override, it will use the provided crosstab
        else:
            _crosstab = crosstab_override.copy()

        _crosstab = _crosstab.replace([np.nan, np.inf, -np.inf],0)
        
        if _crosstab.shape == (0,0):
            _empty_crosstab_error_message = 'The configuration resulted in an empty crosstab. If a custom crosstab was passes make sure it is not an empty frame' 
            raise ValueError(_empty_crosstab_error_message)

        # Load data into category chart data object 
        _chart_data = CategoryChartData()

        for category in _crosstab.index.to_list():
            _chart_data.add_category(category)
        
        for serie in _crosstab.columns.to_list():
            _chart_data.add_series(
                name=serie,
                values=_crosstab[serie].to_list(),
                number_format=number_format_series
            )
        
        # If we are in multiple graph triplet slides mode this will use the slide object for the batch 
        if use_slide_object:
            _graph_slide = slide_object
        else:
            # Initialize the slide if we are creating it for the single slide single graph mode 
            _graph_slide = press_object.slides.add_slide(layout_object)

        if not use_slide_object:
            # Add title to the slide if we are creating it for the single slide single graph mode 
            _graph_slide.placeholders[0].text = slide_title

        # Initialize the graphic frame object 
        _graphic_frame = _graph_slide.shapes.add_chart(
            chart_type=chart_type,
            x=pos_horizontal,
            y=pos_vertical,
            cx=width,
            cy=height,
            chart_data=_chart_data
        )

        # Chart object for decorating
        _chart = _graphic_frame.chart 

        # Adding data lables as % by category
        if show_data_labels_perc:
            # Normalize the crosstab by categories (rows)
            _normalized_crosstab = _crosstab.div(_crosstab.sum(axis=1), axis=0).fillna(0) * 100

            # Insertar etiquetas personalizadas en el XML del gráfico
            _chart_xml = _chart.plots[0].chart._element
            for s_idx, ser in enumerate(_chart_xml.findall('.//c:ser', namespaces=_chart_xml.nsmap)):
                dLbls = ser.find('c:dLbls', namespaces=_chart_xml.nsmap)
                if dLbls is None:
                    dLbls = ET.SubElement(ser, '{http://schemas.openxmlformats.org/drawingml/2006/chart}dLbls')
                for idx, perc in enumerate(_normalized_crosstab.iloc[:, s_idx].tolist()):
                    dLbl = ET.SubElement(dLbls, '{http://schemas.openxmlformats.org/drawingml/2006/chart}dLbl')
                    idx_tag = ET.SubElement(dLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}idx')
                    idx_tag.set('val', str(idx))
                    tx = ET.SubElement(dLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
                    rich = ET.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
                    ET.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                    ET.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
                    p = ET.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
                    r = ET.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
                    ET.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    t = ET.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    t.text = f"{perc:.1f}%"
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showLegendKey', val='0'))
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showVal', val='0'))
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showCatName', val='0'))
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showSerName', val='0'))
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showPercent', val='0'))
                    dLbl.append(ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/chart}showBubbleSize', val='0'))

        # Decoration parameters 
        _chart.has_title = False # Turns off the title for the graph 
        _chart.font.name = 'Arial' # Font for graph
        _chart.font.size = Pt(9) # Font size

        # Turning off gridlines 
        _chart.value_axis.has_major_gridlines = False 
        _chart.value_axis.has_minor_gridlines = False
        _chart.category_axis.has_major_gridlines = False
        _chart.category_axis.has_minor_gridlines = False 

        # Making the value axis start at 0
        if zero_axis_start:
            _chart.value_axis.minimum_scale = 0

        # Turning on or off the category axis 
        if not show_category_axis:
            _chart.category_axis.tick_labels.font.size = Pt(1)
            _chart.category_axis.tick_labels.font.color.rgb = RGBColor(255,255,255)

        # Turining off tickmarks 
        _chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
        _chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE

        # Adding legend 
        _chart.has_legend = show_legend
        if show_legend:
            _chart.legend.position = XL_LEGEND_POSITION.TOP
            _chart.legend.include_in_layout = True

        if show_data_labels_value is not None: 
            for i in range(len(_crosstab.columns.to_list())):
                _chart.series[i].data_labels.show_value = True
                _chart.series[i].data_labels.position = show_data_labels_value # Should be XL_LABEL_POSITION.POSITION kind of object


        # Series color fill 
        _extended_pallete = color_pallete

        while len(_extended_pallete) < len(_crosstab.columns.to_list()):
            _extended_pallete.extend(_extended_pallete)

        for i in range(len(_crosstab.columns.to_list())):
            _chart.series[i].invert_if_negative = False # Disables color inversion for negative values
            _chart.series[i].format.fill.solid()
            _chart.series[i].format.fill.fore_color.rgb = _extended_pallete[i]

            _chart.series[i].format.line.fill.solid()
            _chart.series[i].format.line.fill.fore_color.rgb = _extended_pallete[i]
            _chart.series[i].format.line.width = Pt(1.5)

    elif (len(iterable_vars) > 0) & (len(grouping_vars) == 0):
        
        # Checks if iterable vars columns are all str type 
        check_columns_strtype(_filtered_df, iterable_vars)

        # Creates nested iterable vars 
        _filtered_df['agg_iterable_var'] = _filtered_df[iterable_vars].agg(' '.join, axis=1)

        for _iteration in _filtered_df['agg_iterable_var'].unique():

            # Create the slide for the iteration
            _iteration_slide = press_object.slides.add_slide(layout_object)

            # Add the title of the iteration slide 
            _iteration_slide.placeholders[0].text = ' – '.join([slide_title,_iteration])

            # Adding the iteration chart to the slide 
            chart_maker(
                press_object=press_object,
                layout_object=layout_object,
                use_slide_object=True,
                slide_object=_iteration_slide,
                slide_title='',
                stacked_df=_filtered_df,
                value_var=value_var,
                category_vars=category_vars,
                series_vars=series_vars,
                iterable_vars=[],
                grouping_vars=[],
                filter_mask=_filtered_df['agg_iterable_var']==_iteration,
                chart_type=chart_type,
                show_legend=show_legend,
                show_category_axis=show_category_axis,
                show_data_labels_perc=show_data_labels_perc,
                show_data_labels_value=show_data_labels_value,
                pos_horizontal=pos_horizontal,
                pos_vertical=pos_vertical,
                width=width,
                height=height,
                agg_function=agg_function,
                color_pallete=color_pallete,
                number_format_series=number_format_series,
                zero_axis_start=zero_axis_start
            )


    elif (len(iterable_vars) > 0) & (len(grouping_vars) > 0):

        # Checks if iterable vars columns are all str type 
        check_columns_strtype(_filtered_df, iterable_vars)
        check_columns_strtype(_filtered_df, grouping_vars)

        # Creates a variable with full grouping and iterable vars
        _group_iter_vars = []
        _group_iter_vars.extend(grouping_vars)
        _group_iter_vars.extend(iterable_vars)

        # Creates nested iterable vars 
        _filtered_df['agg_iterable_var'] = _filtered_df[iterable_vars].agg(' '.join, axis=1)
        _filtered_df['agg_grouping_var'] = _filtered_df[grouping_vars].agg(' '.join, axis=1)
        _filtered_df['agg_group_iter_var'] = _filtered_df[_group_iter_vars].agg(' '.join, axis=1)

        for _group in _filtered_df['agg_grouping_var'].unique():

            # Elements that will form the batches we need to make graphs on 
            _items = _filtered_df[_filtered_df['agg_grouping_var']==_group]['agg_group_iter_var'].unique()

            # Batches grouped in elements of 3 
            _batches = [_items[i:i+3] for i in range(0, len(_items), 3)]

            for _batch in _batches:

                # Create the slide for the batch 
                _batch_slide = press_object.slides.add_slide(layout_object)

                # Adds the group as the title 
                _batch_slide.placeholders[0].text = ' – '.join([slide_title,_group])

                for i in range(len(_batch)):

                    # Conditional parameters if it is the last graph on the slide
                    if i + 1 == len(_batch):
                        _height_conditional = Inches(1.86)
                        _show_category_axis_conditional = True

                    else:
                        _height_conditional = Inches(1.43)
                        _show_category_axis_conditional = False

                    if i == 0:
                        # Creates the graph for the first position 
                        chart_maker(
                            press_object=press_object,
                            layout_object=layout_object,
                            use_slide_object=True,
                            slide_object=_batch_slide,
                            slide_title='',
                            stacked_df=_filtered_df,
                            value_var=value_var,
                            category_vars=category_vars,
                            series_vars=series_vars,
                            iterable_vars=[],
                            grouping_vars=[],
                            filter_mask=_filtered_df['agg_group_iter_var'] == _batch[i],
                            chart_type=chart_type,
                            show_legend=show_legend,
                            show_category_axis=_show_category_axis_conditional,
                            show_data_labels_perc=show_data_labels_perc,
                            show_data_labels_value=show_data_labels_value,
                            pos_horizontal=Inches(2.16),
                            pos_vertical=Inches(1.87),
                            width=Inches(10.77),
                            height=_height_conditional,
                            agg_function=agg_function,
                            color_pallete=color_pallete,
                            number_format_series=number_format_series,
                            zero_axis_start=zero_axis_start
                        )

                        # Add textbox (new version)
                        textbox_maker(
                            slide_object=_batch_slide,
                            text=_batch[i],
                            font_type='Arial',
                            font_size=Pt(12),
                            font_color=RGBColor(51,51,51),
                            word_wrap=True,
                            pos_horizontal=Inches(0.39),
                            pos_vertical=Inches(1.87),
                            width=Inches(1.6),
                            height=Inches(0.81)
                        )

                        # # Add textbox
                        # _textbox1 = _batch_slide.shapes.add_textbox(
                        #     Inches(0.39), # Horizontal position
                        #     Inches(1.87), # Vertical position
                        #     Inches(1.6), # Width
                        #     Inches(0.81) # Height
                        # )

                        # # Decorate textbox 
                        # _text_frame1 = _textbox1.text_frame
                        # _text_frame1.text = _batch[i]
                        # _text_frame1.word_wrap = True

                        # # Text format 
                        # _paragraph1 = _text_frame1.paragraphs[0]
                        # _run1 = _paragraph1.runs[0]
                        # _run1.font.name = 'Arial'
                        # _run1.font.size = Pt(12)


                    elif i == 1:
                        # Creates the graph for the second position 
                        chart_maker(
                            press_object=press_object,
                            layout_object=layout_object,
                            use_slide_object=True,
                            slide_object=_batch_slide,
                            slide_title='',
                            stacked_df=_filtered_df,
                            value_var=value_var,
                            category_vars=category_vars,
                            series_vars=series_vars,
                            iterable_vars=[],
                            grouping_vars=[],
                            filter_mask=_filtered_df['agg_group_iter_var'] == _batch[i],
                            chart_type=chart_type,
                            show_legend=show_legend,
                            show_category_axis=_show_category_axis_conditional,
                            show_data_labels_perc=show_data_labels_perc,
                            show_data_labels_value=show_data_labels_value,
                            pos_horizontal=Inches(2.16),
                            pos_vertical=Inches(3.32),
                            width=Inches(10.77),
                            height=_height_conditional,
                            agg_function=agg_function,
                            color_pallete=color_pallete,
                            number_format_series=number_format_series,
                            zero_axis_start=zero_axis_start
                        )

                        # Add textbox (new version)
                        textbox_maker(
                            slide_object=_batch_slide,
                            text=_batch[i],
                            font_type='Arial',
                            font_size=Pt(12),
                            font_color=RGBColor(51,51,51),
                            word_wrap=True,
                            pos_horizontal=Inches(0.39),
                            pos_vertical=Inches(3.32),
                            width=Inches(1.6),
                            height=Inches(0.81)
                        )

                        # # Add textbox
                        # _textbox2 = _batch_slide.shapes.add_textbox(
                        #     Inches(0.39), # Horizontal position
                        #     Inches(3.32), # Vertical position
                        #     Inches(1.6), # Width
                        #     Inches(0.81) # Height
                        # )

                        # # Decorate textbox 
                        # _text_frame2 = _textbox2.text_frame
                        # _text_frame2.text = _batch[i]
                        # _text_frame2.word_wrap = True

                        # # Text format 
                        # _paragraph2 = _text_frame2.paragraphs[0]
                        # _run2 = _paragraph2.runs[0]
                        # _run2.font.name = 'Arial'
                        # _run2.font.size = Pt(12)

                    elif i == 2:
                        # Creates the graph for the third position 

                        chart_maker(
                            press_object=press_object,
                            layout_object=layout_object,
                            use_slide_object=True,
                            slide_object=_batch_slide,
                            slide_title='',
                            stacked_df=_filtered_df,
                            value_var=value_var,
                            category_vars=category_vars,
                            series_vars=series_vars,
                            iterable_vars=[],
                            grouping_vars=[],
                            filter_mask=_filtered_df['agg_group_iter_var'] == _batch[i],
                            chart_type=chart_type,
                            show_legend=show_legend,
                            show_category_axis=_show_category_axis_conditional,
                            show_data_labels_perc=show_data_labels_perc,
                            show_data_labels_value=show_data_labels_value,
                            pos_horizontal=Inches(2.16),
                            pos_vertical=Inches(4.77),
                            width=Inches(10.77),
                            height=_height_conditional,
                            agg_function=agg_function,
                            color_pallete=color_pallete,
                            number_format_series=number_format_series,
                            zero_axis_start=zero_axis_start
                        )

                        # Add textbox (new version)
                        textbox_maker(
                            slide_object=_batch_slide,
                            text=_batch[i],
                            font_type='Arial',
                            font_size=Pt(12),
                            font_color=RGBColor(51,51,51),
                            word_wrap=True,
                            pos_horizontal=Inches(0.39),
                            pos_vertical=Inches(4.77),
                            width=Inches(1.6),
                            height=Inches(0.81)
                        )

                        # # Add textbox
                        # _textbox3 = _batch_slide.shapes.add_textbox(
                        #     Inches(0.39), # Horizontal position
                        #     Inches(4.77), # Vertical position
                        #     Inches(1.6), # Width
                        #     Inches(0.81) # Height
                        # )

                        # # Decorate textbox 
                        # _text_frame3 = _textbox3.text_frame
                        # _text_frame3.text = _batch[i]
                        # _text_frame3.word_wrap = True

                        # # Text format 
                        # _paragraph3 = _text_frame3.paragraphs[0]
                        # _run3 = _paragraph3.runs[0]
                        # _run3.font.name = 'Arial'
                        # _run3.font.size = Pt(12)

    else: 
        _error_message = 'Invalid configuration. grouping_var list was passed but iterable_vars list is empty'
        raise ValueError(_error_message)


# Table maker function
def table_maker(
        press_object:Presentation,
        layout_object,
        use_slide_object:bool=False,
        slide_object=None,
        slide_title:str='title',
        stacked_df:pd.DataFrame=pd.DataFrame(),
        value_var:str='',
        category_vars:list=[],
        series_vars:list=[],
        iterable_vars:list=[],
        grouping_vars:list=[],
        filter_mask = None,
        use_multilevel_categories:bool=False,
        table_value_scale:float=1,
        decimal_places:int=0,
        show_PoP_change:bool=False,
        show_share_perc:bool=False,
        pos_horizontal:Inches=Inches(0.39),
        pos_vertical:Inches=Inches(1.87),
        width:Inches=Inches(12.54),
        height:Inches=Inches(4.37),
        agg_function='sum',
        crosstab_override:pd.DataFrame=None
        ):
    
    """
    Creates table slides with 3 modes of operation:
    
    - Single table in a slide -
    This is the most simple mode, this will create a new slide with a single table on it. For this mode of operation both iterable_vars and grouping vars should be left empty. 

    - Multiple single table slides -
    This mode generates multiple new slides with a single table on each. The table will be created by looping on the unique elements of iterable_vars object. 

    - Multiple table triplet slides - 
    This is the most complex mode, this will create up to 2 tables in a single slide per group of every unique element in the grouping_vars object. 
    If a group would create more than 2 tables it will create more slides as tuples of tables.
    This mode require sboth iterable_vars and grouping_vars lists. 
    """

    if use_multilevel_categories:
        _multilevel_index_sep = ':'
    else:
        _multilevel_index_sep = ' '

    # Filter the data
    if filter_mask is None:
        _filtered_df = stacked_df.copy()
    else:
        _filtered_df = stacked_df[filter_mask].copy()

    # Creation of a single table 
    if (len(iterable_vars) == 0) & (len(grouping_vars) == 0):

        ## Prepare the data, create the crosstab

        # If there is no override, it'll use the specified parameters in the function to make the crosstab
        if crosstab_override is None:
            _crosstab = get_crosstab(
                stacked_df=_filtered_df,
                index_vars=category_vars,
                column_vars=series_vars,
                value_var=value_var,
                agg_function=agg_function,
                str_separator=_multilevel_index_sep
            ).fillna(0)
        
        # If there is a crosstab override, it will use the provided crosstab
        else:
            _crosstab = crosstab_override.copy()

        # Raises error if the crosstab configured or passes is empty
        if _crosstab.shape == (0,0):
            _empty_crosstab_error_message = 'The configuration resulted in an empty crosstab. If a custom crosstab was passes make sure it is not an empty frame' 
            raise ValueError(_empty_crosstab_error_message)

        # Scaling for custom values 
        _crosstab = _crosstab/table_value_scale

        # Creation of the period over period change df 
        if show_PoP_change:
            _crosstab_change = (_crosstab.T.diff(1) / _crosstab.T.shift(1)).T * 100
        
        # Creation of the share df 
        if show_share_perc:
            _crosstab_share = (_crosstab / _crosstab.sum()) * 100

        # Merge with original crosstab
        if show_PoP_change:
            _crosstab = pd.merge(
                left=_crosstab,
                right=_crosstab_change,
                left_index=True,
                right_index=True,
                suffixes=(None,' change')
            )
        
        if show_share_perc:
            _crosstab = pd.merge(
                left=_crosstab,
                right=_crosstab_share,
                left_index=True,
                right_index=True,
                suffixes=(None,' share')
            )
        
        # Number rounding 
        _crosstab = _crosstab.round(decimal_places)

        # Creates the multilevel index if requested
        if use_multilevel_categories:
            _crosstab.index = pd.MultiIndex.from_tuples(_crosstab.index.str.split(_multilevel_index_sep).map(tuple))

            # If we are in multiple graph triplet slides mode this will use the slide object for the batch 
        if use_slide_object:
            _table_slide = slide_object
        else:
            # Initialize the slide if we are creating it for the single slide single graph mode 
            _table_slide = press_object.slides.add_slide(layout_object)

        if not use_slide_object:
            # Add title to the slide if we are creating it for the single slide single graph mode 
            _table_slide.placeholders[0].text = slide_title

        
        ## Creating the table object 

        # Content shape 
        _rows, _cols = _crosstab.shape

        if use_multilevel_categories:
            _index_levels = len(_crosstab.index.to_list()[0])
        else:
            _index_levels = 1

        _table_shape = _table_slide.shapes.add_table(
            _rows + 1, 
            _cols + _index_levels,
            pos_horizontal,
            pos_vertical,
            width,
            height
        ).table

        # Filling the table headers (columns) 
        for _col_idx, _col_name in enumerate(_crosstab.columns):
            _cell = _table_shape.cell(0, _col_idx + _index_levels)
            _cell.text = str(_col_name)
            _cell.fill.solid()
            _cell.fill.fore_color.rgb = RGBColor(51,51,51)
            _cell.text_frame.paragraphs[0].font.bold = True
            _cell.text_frame.paragraphs[0].font.name = 'Arial'
            _cell.text_frame.paragraphs[0].font.size = Pt(10.5)
            _cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Formating the rest of the column cells 
        for _col_idx in range(_index_levels):
            _cell = _table_shape.cell(0, _col_idx)
            _cell.fill.solid()
            _cell.fill.fore_color.rgb = RGBColor(51,51,51)
            _cell.text_frame.paragraphs[0].font.bold = True
            _cell.text_frame.paragraphs[0].font.name = 'Arial'
            _cell.text_frame.paragraphs[0].font.size = Pt(10.5)
            _cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Filling the table categories (index / multilevel index)
        for _index_item_idx in range(_rows):
            for _index_level_idx in range(_index_levels):
                _cell = _table_shape.cell(_index_item_idx + 1, _index_level_idx)
                if use_multilevel_categories:
                    _cell.text = str(_crosstab.index[_index_item_idx][_index_level_idx])
                else:
                    _cell.text = str(_crosstab.index[_index_item_idx])
                _cell.fill.solid()
                _cell.fill.fore_color.rgb = RGBColor(255,255,255)
                _cell.text_frame.paragraphs[0].font.bold = False
                _cell.text_frame.paragraphs[0].font.name = 'Arial'
                _cell.text_frame.paragraphs[0].font.size = Pt(9)
                _cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT               

        # Filling the table values 
        for _row_idx in range(_rows):
            for _col_idx in range(_cols):
                _cell = _table_shape.cell(_row_idx + 1, _col_idx + _index_levels)
                _value = _crosstab.iloc[_row_idx, _col_idx]
                if isinstance(_value, (int, float)):
                    _formatted = _formatted = f'{_value:,.{decimal_places}f}' if decimal_places > 0 else f'{_value:,.0f}'
                    _cell.text = _formatted
                else:
                    _cell.text = str(_value)
                _cell.fill.solid()
                _cell.fill.fore_color.rgb = RGBColor(255,255,255)
                _cell.text_frame.paragraphs[0].font.bold = False
                _cell.text_frame.paragraphs[0].font.name = 'Arial'
                _cell.text_frame.paragraphs[0].font.size = Pt(9)
                _cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT 

        # Merging cells in the index if they are the same value 
        for _index_level_idx in range(_index_levels):

            # Start the objects that keep track of the ranges that need merging
            _start_row_idx = 0
            _end_row_idx = None

            for _row_idx in range(_rows):
                # If it's the last element in the rows 
                if _row_idx == _rows - 1:
                    _end_row_idx = _row_idx
                    # Do the last merge with the last starting cell
                    # First save the text of the first cell
                    _start_cell_text = _table_shape.cell(_start_row_idx + 1, _index_level_idx).text
                    # Doing the merge
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).merge(
                        _table_shape.cell(_end_row_idx + 1, _index_level_idx)
                        )
                    # Updating the text to only have the first one
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).text = _start_cell_text

                    # Formating 
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).fill.fore_color.rgb = RGBColor(255,255,255)
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.bold = False
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.name = 'Arial'
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.size = Pt(9)
                    _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.alignment = PP_ALIGN.LEFT 


                        
                else:
                    if use_multilevel_categories:
                        # If the next element is different from the current one
                        if _crosstab.index[_row_idx][_index_level_idx] != _crosstab.index[_row_idx + 1][_index_level_idx]:
                            # Flag the current element as the end point 
                            _end_row_idx = _row_idx
                            # Merge the cells
                            # First save the text of the first cell
                            _start_cell_text = _table_shape.cell(_start_row_idx + 1, _index_level_idx).text
                            # Do the merge
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).merge(
                                _table_shape.cell(_end_row_idx + 1, _index_level_idx)
                                )
                            # Updating the text to only have the first one
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).text = _start_cell_text

                            # Formating 
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).fill.fore_color.rgb = RGBColor(255,255,255)
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.bold = False
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.name = 'Arial'
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.size = Pt(9)
                            _table_shape.cell(_start_row_idx + 1, _index_level_idx).text_frame.paragraphs[0].font.alignment = PP_ALIGN.LEFT 

                            # Update the start of the next range
                            _start_row_idx = _row_idx + 1
                    
                    elif not use_multilevel_categories:
                        # If the next element is different from the current one
                        if _crosstab.index[_row_idx] != _crosstab.index[_row_idx + 1]:
                            # Flag the current element as the end point 
                            _end_row_idx = _row_idx
                            # Merge the cells
                            # First save the text of the first cell
                            _start_cell_text = _table_shape.cell(_start_row_idx + 1, 0).text
                            # Do the merge
                            _table_shape.cell(_start_row_idx + 1, 0).merge(
                                _table_shape.cell(_end_row_idx + 1, 0)
                                )
                            # Updating the text to only have the first one
                            _table_shape.cell(_start_row_idx + 1, 0).text = _start_cell_text

                            # Formating 
                            _table_shape.cell(_start_row_idx + 1, 0).fill.fore_color.rgb = RGBColor(255,255,255)
                            _table_shape.cell(_start_row_idx + 1, 0).text_frame.paragraphs[0].font.bold = False
                            _table_shape.cell(_start_row_idx + 1, 0).text_frame.paragraphs[0].font.name = 'Arial'
                            _table_shape.cell(_start_row_idx + 1, 0).text_frame.paragraphs[0].font.size = Pt(9)
                            _table_shape.cell(_start_row_idx + 1, 0).text_frame.paragraphs[0].font.alignment = PP_ALIGN.LEFT 

                            # Update the start of the next range
                            _start_row_idx = _row_idx + 1

        # Adding borders to the cells 
        for _row_idx in range(_rows):
            for _col_idx in range(_cols + _index_levels):
                set_cell_border(_table_shape.cell(_row_idx+1, _col_idx), border_color_TB='333333', border_color_LR='FFFFFF')
        
        # Adjusting the row height 
        for _row_idx, _row in enumerate(_table_shape.rows):
            _row.height = Inches(0.25)


    # Creation of a series of single table per slide for each iterable element
    elif (len(iterable_vars) > 0) & (len(grouping_vars) == 0):
        
        # Checks if iterable vars columns are all str type 
        check_columns_strtype(_filtered_df, iterable_vars)

        # Creates nested iterable vars 
        _filtered_df['agg_iterable_var'] = _filtered_df[iterable_vars].agg(' '.join, axis=1)

        # Apply the graphing loop again for each iteration. Filters the filtered df by the corresponding iteration and adds as a subtitle the iteration. 
        for _iteration in _filtered_df['agg_iterable_var'].unique():

            # Create the slide for the iteration
            _iteration_slide = press_object.slides.add_slide(layout_object)

            # Add the title of the iteration slide 
            _iteration_slide.placeholders[0].text = ' – '.join([slide_title,_iteration])

            # Adding the iteration table to the slide 
            table_maker(
                press_object=press_object,
                layout_object=layout_object,
                use_slide_object=True,
                slide_object=_iteration_slide,
                slide_title='',
                stacked_df=_filtered_df,
                value_var=value_var,
                category_vars=category_vars,
                series_vars=series_vars,
                iterable_vars=[],
                grouping_vars=[],
                filter_mask=_filtered_df['agg_iterable_var']==_iteration,
                use_multilevel_categories=use_multilevel_categories,
                table_value_scale=table_value_scale,
                decimal_places=decimal_places,
                show_PoP_change=show_PoP_change,
                show_share_perc=show_share_perc,
                pos_horizontal=pos_horizontal,
                pos_vertical=pos_vertical,
                width=width,
                height=height,
                agg_function=agg_function
            )            

    # Creation of multiple tables per slide for each grouping element and iterable element 
    elif (len(iterable_vars) > 0) & (len(grouping_vars) > 0):

        # Checks if iterable vars columns are all str type 
        check_columns_strtype(_filtered_df, iterable_vars)
        check_columns_strtype(_filtered_df, grouping_vars)

        # Creates a variable with full grouping and iterable vars
        _group_iter_vars = []
        _group_iter_vars.extend(grouping_vars)
        _group_iter_vars.extend(iterable_vars)

        # Creates nested iterable vars 
        _filtered_df['agg_iterable_var'] = _filtered_df[iterable_vars].agg(' '.join, axis=1)
        _filtered_df['agg_grouping_var'] = _filtered_df[grouping_vars].agg(' '.join, axis=1)
        _filtered_df['agg_group_iter_var'] = _filtered_df[_group_iter_vars].agg(' '.join, axis=1)

        for _group in _filtered_df['agg_grouping_var'].unique():

            # Elements that will form the batches we need to make graphs on 
            _items = _filtered_df[_filtered_df['agg_grouping_var']==_group]['agg_group_iter_var'].unique()

            # Batches grouped in elements of 2
            _batches = [_items[i:i+2] for i in range(0, len(_items), 2)]

            for _batch in _batches:

                # Create the slide for the batch 
                _batch_slide = press_object.slides.add_slide(layout_object)

                # Adds the group as the title 
                _batch_slide.placeholders[0].text = ' – '.join([slide_title,_group])

                for _batch_idx in range(len(_batch)):

                    if _batch_idx == 0:
                        # Creates the table for the left position
                        table_maker(
                            press_object=press_object,
                            layout_object=layout_object,
                            use_slide_object=True,
                            slide_object=_batch_slide,
                            slide_title='',
                            stacked_df=_filtered_df,
                            value_var=value_var,
                            category_vars=category_vars,
                            series_vars=series_vars,
                            iterable_vars=[],
                            grouping_vars=[],
                            filter_mask=_filtered_df['agg_group_iter_var']==_batch[_batch_idx],
                            use_multilevel_categories=use_multilevel_categories,
                            table_value_scale=table_value_scale,
                            decimal_places=decimal_places,
                            show_PoP_change=show_PoP_change,
                            show_share_perc=show_share_perc,
                            pos_horizontal=Inches(0.39),
                            pos_vertical=Inches(1.87),
                            width=Inches(6.15),
                            height=Inches(4.37),
                            agg_function=agg_function
                        )


                    if _batch_idx == 1:
                        # Creates the table for the right position
                        table_maker(
                            press_object=press_object,
                            layout_object=layout_object,
                            use_slide_object=True,
                            slide_object=_batch_slide,
                            slide_title='',
                            stacked_df=_filtered_df,
                            value_var=value_var,
                            category_vars=category_vars,
                            series_vars=series_vars,
                            iterable_vars=[],
                            grouping_vars=[],
                            filter_mask=_filtered_df['agg_group_iter_var']==_batch[_batch_idx],
                            use_multilevel_categories=use_multilevel_categories,
                            table_value_scale=table_value_scale,
                            decimal_places=decimal_places,
                            show_PoP_change=show_PoP_change,
                            show_share_perc=show_share_perc,
                            pos_horizontal=Inches(6.78),
                            pos_vertical=Inches(1.87),
                            width=Inches(6.15),
                            height=Inches(4.37),
                            agg_function=agg_function
                        )


# Cell border decoration function (DOESN'T WORK WELL)
def set_cell_border(cell, border_color_TB='FF0000', border_color_LR='0000FF', border_width="12700"):
    """
    Applies different color border for a cell. 
    border_color_TB: color for top y bottom (eg. 'FF0000' red)
    border_color_LR: color for left y right (eg. '0000FF' blue)
    border_width: width of border in EMUs (12700 = 1pt)
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    def create_border_line(color):
        ln = OxmlElement("a:ln")
        ln.set("w", border_width)

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", color)
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", "solid")
        ln.append(prstDash)

        ln.append(OxmlElement("a:round"))

        headEnd = OxmlElement("a:headEnd")
        headEnd.set("type", "none")
        ln.append(headEnd)

        tailEnd = OxmlElement("a:tailEnd")
        tailEnd.set("type", "none")
        ln.append(tailEnd)

        return ln

    # Crear y asignar bordes
    for side, color in {
        "lnT": border_color_TB,
        "lnB": border_color_TB,
        "lnL": border_color_LR,
        "lnR": border_color_LR
    }.items():
        ln = create_border_line(color)
        ln.tag = qn(f'a:{side}')
        tcPr.append(ln)

# Scatter plot maker function 
def scatter_maker(
    press_object: Presentation,
    layout_object,
    use_slide_object: bool = False,
    slide_object=None,
    slide_title: str = 'title',
    # Data
    crosstab_df: pd.DataFrame = pd.DataFrame(),
    x_col: str = '',
    y_col: str = '',
    size_col: str | None = None,  # required only for bubble
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.XY_SCATTER,
    # Visuals & position
    pos_horizontal: Inches = Inches(0.39),
    pos_vertical: Inches = Inches(1.87),
    width: Inches = Inches(12.54),
    height: Inches = Inches(4.37),
    color_pallete: list = standard_pallete,
    marker_style = XL_MARKER_STYLE.CIRCLE,
    marker_size: int = 7,
    # Axis behavior
    axis_center_mode: str = 'zero'  # 'zero' | 'median' | 'mean' | 'none'
):
    """
    Create a single XY scatter or Bubble chart slide, styled like chart_maker:
      - accepts a crosstab-like df where the *index* is the label to show per point
      - requires columns: x_col, y_col (and optional size_col for bubbles)
      - hides gridlines, hides tick labels, removes tick marks
      - can center axes at 0 / median / mean of x and y
    Important: This function doesn't allow for the creating of the crosstab within itself, it should be provided a ready to use crosstab without NAs
    """
    # --- Validate configuration
    if chart_type not in (XL_CHART_TYPE.XY_SCATTER, XL_CHART_TYPE.BUBBLE):
        raise ValueError("chart_type must be XL_CHART_TYPE.XY_SCATTER or XL_CHART_TYPE.BUBBLE")

    _required = [x_col, y_col] if chart_type == XL_CHART_TYPE.XY_SCATTER else [x_col, y_col, size_col]
    if any(col is None or col == '' for col in _required):
        raise ValueError("Please provide x_col and y_col (and size_col for bubble).")

    _missing = [c for c in _required if c not in crosstab_df.columns]
    if _missing:
        raise ValueError(f"Missing required column(s) in crosstab_df: {_missing}")

    # Keep only needed columns, coerce to numeric and drop missing rows
    _filtered_df = crosstab_df[[x_col, y_col] + ([size_col] if chart_type == XL_CHART_TYPE.BUBBLE else [])].copy()
    for col in [x_col, y_col] + ([size_col] if chart_type == XL_CHART_TYPE.BUBBLE else []):
        _filtered_df[col] = pd.to_numeric(_filtered_df[col], errors='coerce')
    _filtered_df = _filtered_df.dropna(subset=[x_col, y_col] + ([size_col] if chart_type == XL_CHART_TYPE.BUBBLE else []))
    if _filtered_df.empty:
        raise ValueError("No valid rows after coercing to numeric. Check your crosstab values.")

    # --- Build chart data
    if chart_type == XL_CHART_TYPE.XY_SCATTER:
        _chart_data = XyChartData()
        _series = _chart_data.add_series('Series')
        for _, row in _filtered_df.iterrows():
            _series.add_data_point(float(row[x_col]), float(row[y_col]))
    else:
        _chart_data = BubbleChartData()
        _series = _chart_data.add_series('Series')
        for _, row in _filtered_df.iterrows():
            _series.add_data_point(float(row[x_col]), float(row[y_col]), float(row[size_col]))

    # --- Create or reuse slide
    if use_slide_object and slide_object is not None:
        _graph_slide = slide_object
    else:
        _graph_slide = press_object.slides.add_slide(layout_object)
    
    if not use_slide_object:
        # Add title to the slide if we are creating it for the single slide single graph mode 
        _graph_slide.placeholders[0].text = slide_title

    # Initialize the graphic frame object 
    _graphic_frame = _graph_slide.shapes.add_chart(
        chart_type=chart_type,
        x=pos_horizontal,
        y=pos_vertical,
        cx=width,
        cy=height,
        chart_data=_chart_data
    )

    # Chart object for decorating
    _chart = _graphic_frame.chart

    # Style: fonts, gridlines, ticks, legend
    _chart.has_title = False
    _chart.has_legend = False

    # Resolve axes (for scatter/bubble, category_axis is the X value axis; value_axis is Y)
    _x_axis = _chart.category_axis
    _y_axis = _chart.value_axis

    # No gridlines, no tick marks, no numeric labels (hide by formatting like in chart_maker)
    for _ax in (_x_axis, _y_axis):
        _ax.has_major_gridlines = False
        _ax.has_minor_gridlines = False
        _ax.major_tick_mark = XL_TICK_MARK.NONE
        _ax.minor_tick_mark = XL_TICK_MARK.NONE
        _ax.tick_labels.font.size = Pt(1)
        _ax.tick_labels.font.color.rgb = RGBColor(255, 255, 255)


    # Decoration for XY scatter 
    if chart_type == XL_CHART_TYPE.XY_SCATTER:
        _chart.series[0].marker.style = marker_style
        _chart.series[0].marker.size = marker_size
        _chart.series[0].marker.format.fill.solid()
        _chart.series[0].marker.format.fill.fore_color.rgb = color_pallete[0]
        _chart.series[0].marker.format.line.fill.solid()
        _chart.series[0].marker.format.line.fill.fore_color.rgb = color_pallete[0]

    # Series decoration for Bubble chart
    if chart_type == XL_CHART_TYPE.BUBBLE:
        _chart.series[0].format.fill.solid()
        _chart.series[0].format.fill.fore_color.rgb = color_pallete[0]
        _chart.series[0].format.line.fill.solid()
        _chart.series[0].format.line.fill.fore_color.rgb = color_pallete[0]


    # Axes ranges with padding
    _xs = _filtered_df[x_col].astype(float)
    _ys = _filtered_df[y_col].astype(float)
    def _pad_range(vals: pd.Series, ratio: float):
        vmin, vmax = vals.min(), vals.max()
        span = vmax - vmin
        if span == 0:
            span = max(abs(vmax), 1.0) * 0.1  # degenerate case
        pad = span * ratio
        return (vmin - pad, vmax + pad)

    _x_min, _x_max = _pad_range(_xs, 0.05)
    _y_min, _y_max = _pad_range(_ys, 0.05)

    # Axis centering
    if axis_center_mode in ('zero', 'median', 'mean'):
        _x_center = 0.0 if axis_center_mode == 'zero' else (_xs.median() if axis_center_mode == 'median' else _xs.mean())
        _y_center = 0.0 if axis_center_mode == 'zero' else (_ys.median() if axis_center_mode == 'median' else _ys.mean())
        # Ensure center lies within scale
        _x_min = min(_x_min, _x_center)
        _x_max = max(_x_max, _x_center)
        _y_min = min(_y_min, _y_center)
        _y_max = max(_y_max, _y_center)

        # Where the Y axis crosses the X axis (on X scale)
        _x_axis.crosses = XL_AXIS_CROSSES.CUSTOM
        _x_axis.crosses_at = float(_x_center)
        # Where the X axis crosses the Y axis (on Y scale)
        _y_axis.crosses = XL_AXIS_CROSSES.CUSTOM
        _y_axis.crosses_at = float(_y_center)

    # Apply ranges (after ensuring centers are in-range)
    _x_axis.minimum_scale = float(_x_min)
    _x_axis.maximum_scale = float(_x_max)
    _y_axis.minimum_scale = float(_y_min)
    _y_axis.maximum_scale = float(_y_max)


    # Adding data labels to the charts 
    chart_xml = _chart.plots[0].chart._element
    ser = chart_xml.find('.//c:ser', namespaces=chart_xml.nsmap)

    # Ensure dLbls exists
    dLbls = ser.find('c:dLbls', namespaces=chart_xml.nsmap)
    if dLbls is None:
        dLbls = OxmlElement('c:dLbls')
        ser.append(dLbls)

    for idx, label in enumerate(_filtered_df.index):
        dLbl = OxmlElement('c:dLbl')
        
        idx_tag = OxmlElement('c:idx')
        idx_tag.set('val', str(idx))
        dLbl.append(idx_tag)

        tx = OxmlElement('c:tx')
        rich = OxmlElement('c:rich')
        bodyPr = OxmlElement('a:bodyPr')
        lstStyle = OxmlElement('a:lstStyle')
        p = OxmlElement('a:p')
        r = OxmlElement('a:r')
        rPr = OxmlElement('a:rPr')
        rPr.set('sz', '900')  # 9pt
        t = OxmlElement('a:t')
        t.text = str(label)

        r.append(rPr)
        r.append(t)
        p.append(r)
        rich.append(bodyPr)
        rich.append(lstStyle)
        rich.append(p)
        tx.append(rich)
        dLbl.append(tx)
        
        showLegendKey = OxmlElement('c:showLegendKey')
        showLegendKey.set('val', '0')
        dLbl.append(showLegendKey)


        dLbls.append(dLbl)

# Simplified label creator
def textbox_maker(
        slide_object=None,
        text:str='',
        font_type:str='Arial',
        font_size:Pt=Pt(12),
        font_color:RGBColor=RGBColor(51,51,51), # Black 
        word_wrap:bool=True,
        pos_horizontal:Inches=Inches(0.39),
        pos_vertical:Inches=Inches(1.87),
        width:Inches=Inches(1.25),
        height:Inches=Inches(1.25)
        ):

    """
    This function adds a custom textbox to the slide object selected.
    Font type, size and color can be adjusted. 
    """

    # Check if there is slide object 
    if slide_object is None:
        _error_message = 'There was no slide object passed to the function. Make sure to include it.'
        raise ValueError(_error_message)

    # Create the textbox
    _textbox = slide_object.shapes.add_textbox(
        pos_horizontal,
        pos_vertical,
        width,
        height
    )

    _text_frame = _textbox.text_frame
    _text_frame.clear()  # Clear any default content
    _paragraph = _text_frame.add_paragraph()
    _run = _paragraph.add_run()
    _run.text = text
    _run.font.name = font_type
    _run.font.size = font_size
    _run.font.color.rgb = font_color
